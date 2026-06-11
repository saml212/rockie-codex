#!/usr/bin/env python3
"""render_deck.py -- A2: render findings.json into a partner-grade slide deck.

A2 consumes the typed A1 contract (findings.json, sections 0-8) and produces a
slide deck. It does NOT re-run A1 (ingest/reconcile/research/findings) and it
does NOT run the A3 senior-partner critic loop.

Two output paths, by design:

  [primary] --emit payload   Map findings -> the `powerpoint` skill request
                             payload (the canonical Rockie deck engine:
                             rockie-pptagent-render + emit_artifact kind=slides,
                             30-slide cap). This is a pure data transform and is
                             always runnable -- it is what A2 hands to the
                             runtime powerpoint tool.

  [fallback] --emit pptx     Render a real .pptx LOCALLY with python-pptx if it
                             is importable (no network, no GPU, no DeepPresenter).
             --emit json     If python-pptx is absent, dump the deck spec as JSON
                             so the findings->slides mapping is still inspectable.
             --emit auto      pptx if python-pptx present, else json (default).

  [A5 OUT]   --emit emit-args --pptx-in <deck.pptx>
                             Build the exact `emit_artifact` arguments body (kind
                             =slides, base64 content, .pptx filename, notebook_id
                             =$PLATFORM_LAB_ID, destinations=[ui,chat]) that ships
                             the rendered deck OUT to the user as a downloadable
                             artifact. The agent POSTs this to
                             /api/agent-tools/emit_artifact (or calls the
                             emit_artifact MCP tool).

Slide order (pyramid principle):
  1  Title / executive summary  (Section 0 FIRST -- the verdict the IC reads)
  2..N  one slide per Section 1..8: declarative `headline` as the slide title,
        key_facts as bullets that PRESERVE the citation [file | "quote"] and
        the confidence tag. "Not available" stays explicit.
  +1  consolidated Risk register (H/M/L across all sections)
  +1  Open questions for management (pooled across all sections)

Evidence discipline: every fact bullet carries its citation. A fact with no
verbatim quote is rendered as an explicit "[no citation]" marker, never silently
dropped -- so a missing citation is visible to the partner, not hidden.

Respects the powerpoint skill's hard 30-slide cap and A2's 18-24 target band.

Usage:
  render_deck.py <findings.json> [--emit auto|payload|pptx|json]
                 [--out PATH] [--title "Deck title"] [--company NAME]
  Dependencies: python3 stdlib. python-pptx used ONLY for --emit pptx.
"""

import sys
import os
import json
import argparse
import base64
import tempfile

# ----- constants ------------------------------------------------------------

HARD_SLIDE_CAP = 30          # the powerpoint skill rejects > 30 slides
TARGET_MIN, TARGET_MAX = 18, 24  # A2 target band (informational warning only)
CONF_TAG = {"high": "H", "medium": "M", "low": "L"}
RATING_ORDER = {"H": 0, "M": 1, "L": 2}

# A5 deck-OUT: the office MIME the powerpoint skill emits for kind=slides.
PPTX_MIME = ("application/vnd.openxmlformats-officedocument"
             ".presentationml.presentation")


# ----- load + validate ------------------------------------------------------

def load_findings(path):
    try:
        with open(path, "r", encoding="utf-8") as fh:
            data = json.load(fh)
    except FileNotFoundError:
        die("findings file not found: %s" % path)
    except json.JSONDecodeError as e:
        die("findings.json is not valid JSON (%s): line %d column %d"
            % (path, e.lineno, e.colno))
    if not isinstance(data, list):
        die("findings.json must be a JSON array of section objects")
    by_id = {}
    for sec in data:
        sid = sec.get("section_id")
        if sid is None:
            die("a section is missing section_id")
        by_id[sid] = sec
    if 0 not in by_id:
        die("Section 0 (executive summary) is REQUIRED and must lead the deck")
    return data, by_id


def die(msg):
    sys.stderr.write("render_deck: ERROR: %s\n" % msg)
    sys.exit(1)


def default_out(suffix):
    """Pick a default output path in the system temp dir, NOT the cwd.

    Without --out the renderer must never drop build artifacts
    (deck_payload.json / deck.pptx / deck_spec.json) into whatever
    directory it happens to run in -- that pollutes the repo. A temp
    file is created and its path returned.
    """
    fd, path = tempfile.mkstemp(prefix="render_deck_", suffix=suffix)
    os.close(fd)
    return path


# ----- citation formatting (evidence discipline) ----------------------------

def cite_suffix(fact):
    """Build the [file @ locator | "quote"] (conf) suffix.

    Never silently drop a cite -- and never drop the cite_locator (p.N,
    "Section 4.3", URL, ...): a locator is what lets a partner verify the
    quote against the source, so it must survive into the rendered bullet.
    """
    cf = (fact.get("cite_file") or "").strip()
    cq = (fact.get("cite_quote") or "").strip()
    cl = (fact.get("cite_locator") or "").strip()
    conf = CONF_TAG.get((fact.get("confidence") or "").lower(), "?")
    loc = "" if cl in ("", "n/a") else " @ " + cl
    if not cq:
        # missing citation must be VISIBLE, not hidden -- partner-grade rule
        return "  [no citation%s] (%s)" % (loc, conf)
    if cf in ("", "none"):
        # "Not available" facts carry the explicit marker in cite_quote
        head = cl + " | " if loc else ""
        return '  [%s"%s"] (%s)' % (head, cq, conf)
    return '  [%s%s | "%s"] (%s)' % (cf, loc, cq, conf)


def fact_bullet(fact):
    return (fact.get("fact") or "(empty fact)").strip() + cite_suffix(fact)


def fact_citation(fact):
    """Structured citation for a fact -- carries cite_locator through to the
    payload so the downstream engine has the full [file @ locator | quote]
    triple, not just the flattened bullet string."""
    return {
        "fact": (fact.get("fact") or "").strip(),
        "cite_file": (fact.get("cite_file") or "").strip(),
        "cite_quote": (fact.get("cite_quote") or "").strip(),
        "cite_locator": (fact.get("cite_locator") or "").strip(),
        "confidence": (fact.get("confidence") or "").strip(),
    }


# ----- findings -> ordered slide spec ---------------------------------------

def build_slides(findings, by_id, title, company):
    slides = []

    # Slide 1: Section 0 FIRST -- executive summary & recommendation
    s0 = by_id[0]
    rec = (s0.get("recommendation") or "").strip()
    rec_label = {
        "proceed": "PROCEED",
        "proceed_with_conditions": "PROCEED WITH CONDITIONS",
        "pass": "PASS",
    }.get(rec, rec.upper() or "(recommendation not set)")
    exec_bullets = [fact_bullet(f) for f in s0.get("key_facts", [])]
    slides.append({
        "layout": "title",
        "title": s0.get("headline") or "Executive summary & recommendation",
        "subtitle": "%s  |  Recommendation: %s" % (
            company or "Acquisition due diligence", rec_label),
        "bullets": exec_bullets,
        "facts": [fact_citation(f) for f in s0.get("key_facts", [])],
        "section_id": 0,
    })

    # Slides 2..N: one per section with integer section_id >= 1, in order.
    # Iterate EVERY such section (not a hardcoded 1..8 window): sections 9+
    # must not be silently dropped, and only this makes the 30-slide cap
    # reachable for a large findings set.
    body_ids = sorted(
        sid for sid in by_id
        if isinstance(sid, int) and not isinstance(sid, bool) and sid >= 1)
    for sid in body_ids:
        sec = by_id[sid]
        bullets = [fact_bullet(f) for f in sec.get("key_facts", [])]
        # surface reconcile flags inline so a cross-doc delta is never lost
        for rf in sec.get("reconcile_flags", []) or []:
            bullets.append("RECONCILE: " + rf)
        slides.append({
            "layout": "content",
            "title": sec.get("headline") or sec.get("section") or ("Section %d" % sid),
            "bullets": bullets,
            "facts": [fact_citation(f) for f in sec.get("key_facts", [])],
            "section_id": sid,
        })

    # Consolidated risk register (H/M/L), pooled across all sections
    risks = []
    for sec in findings:
        for r in sec.get("risks", []) or []:
            risks.append((r.get("rating", "?"), r.get("risk", ""),
                          sec.get("section", "")))
    risks.sort(key=lambda t: RATING_ORDER.get(t[0], 9))
    if risks:
        slides.append({
            "layout": "content",
            "title": "Risk register: %d flagged (%dH / %dM / %dL)" % (
                len(risks),
                sum(1 for r in risks if r[0] == "H"),
                sum(1 for r in risks if r[0] == "M"),
                sum(1 for r in risks if r[0] == "L")),
            "bullets": ["[%s] %s  (%s)" % (rt, rk, sn) for rt, rk, sn in risks],
            "section_id": "risks",
        })

    # Open questions for management, pooled across all sections
    questions = []
    for sec in findings:
        for q in sec.get("open_questions", []) or []:
            questions.append(q)
    if questions:
        slides.append({
            "layout": "content",
            "title": "Open questions for management (%d)" % len(questions),
            "bullets": questions,
            "section_id": "open_questions",
        })

    return slides


# ----- emit: powerpoint-skill request payload (canonical integration) -------

def to_pptagent_payload(slides, title):
    """Map the slide spec to the `powerpoint` skill's documented Request Shape.

    The contract (../powerpoint/SKILL.md) accepts EXACTLY these top-level keys:

        {prompt, slide_count, title, findings[], attachments[],
         template_path, theme{}}

    It does NOT accept an invented `slides[]` array, nor a `template`
    string. The per-slide structure therefore rides inside the documented
    `findings[]` array (whose item schema the contract leaves open), one
    entry per slide. Each entry carries its structured citations -- including
    cite_locator -- so the downstream engine keeps the full
    [file @ locator | quote] triple rather than a flattened bullet string.
    """
    if len(slides) > HARD_SLIDE_CAP:
        die("deck has %d slides; powerpoint skill caps at %d -- trim sections"
            % (len(slides), HARD_SLIDE_CAP))

    findings = []
    for s in slides:
        entry = {
            "section_id": s.get("section_id"),
            "layout": s["layout"],
            "headline": s["title"],
            "bullets": s.get("bullets", []),
        }
        if s.get("subtitle"):
            entry["subtitle"] = s["subtitle"]
        if s.get("facts"):
            # structured citations carry cite_locator through (not just bullets)
            entry["facts"] = s["facts"]
        findings.append(entry)

    prompt = (
        "Render a partner-grade acquisition due-diligence deck titled "
        "%r. Section 0 (executive summary & recommendation) MUST lead. "
        "Use one slide per entry in `findings`, preserving each fact's "
        "verbatim citation [cite_file @ cite_locator | cite_quote] and "
        "confidence tag. Do not drop citation locators." % title
    )

    # Conform to the documented Request Shape EXACTLY -- no extra keys.
    return {
        "prompt": prompt,
        "slide_count": len(findings),
        "title": title,
        "findings": findings,
        "attachments": [],
        "template_path": None,        # use the bundled Pebble ML template
        "theme": {},                  # engine default theme
    }


# ----- emit: ready-to-post emit_artifact arguments (A5 deck-OUT) -------------

def build_emit_args(pptx_path, title, slide_count, notebook_id, filename=None):
    """Build the EXACT `emit_artifact` arguments body the agent posts to ship
    the rendered deck OUT to the user as a downloadable artifact.

    The deck leaves the runtime via the agent-tools `emit_artifact` tool
    (platform-context/api/agent_tools/schemas.py:597, kind enum includes
    "slides"; handler emit_artifact in api/services/artifact_emitter.py).
    For office artifacts the tool REQUIRES base64-encoded bytes with
    content_encoding="base64" and a .pptx filename; required keys are
    {kind, content, title, notebook_id}. destinations default to the
    powerpoint skill's ["ui", "chat"], the two always-available channels —
    "ui" makes it downloadable in the lab; "chat" drops it into the thread.

    notebook_id is the LAB id (PLATFORM_LAB_ID in the runtime). The mcp-rockie
    bridge sends X-Tenant-Token / X-Tenant-Id so the artifact is scoped to the
    tenant; the artifact then lists/streams via /api/artifacts/{id}/file.
    """
    with open(pptx_path, "rb") as fh:
        raw = fh.read()
    b64 = base64.b64encode(raw).decode("ascii")
    fname = filename or "diligence-deck.pptx"
    if not fname.endswith(".pptx"):
        fname += ".pptx"
    return {
        "kind": "slides",
        "content": b64,                 # raw base64, no "base64:" prefix
        "content_encoding": "base64",
        "filename": fname,
        "mime_type": PPTX_MIME,
        "title": title,
        "notebook_id": notebook_id,
        "destinations": ["ui", "chat"],
        "metadata": {"skill": "diligence-deck", "slide_count": slide_count},
    }


# ----- emit: local .pptx via python-pptx (dependency-light fallback) ---------

def render_pptx(slides, out_path, title):
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except Exception:
        return False  # caller falls back to JSON
    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]
    for s in slides:
        if s["layout"] == "title":
            slide = prs.slides.add_slide(title_layout)
            slide.shapes.title.text = s["title"]
            if slide.placeholders and len(slide.placeholders) > 1:
                slide.placeholders[1].text = s.get("subtitle", "")
            # exec bullets in a textbox below the subtitle
            if s.get("bullets"):
                box = slide.shapes.add_textbox(Pt(40), Pt(200), Pt(640), Pt(300))
                tf = box.text_frame
                tf.word_wrap = True
                for i, b in enumerate(s["bullets"]):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = "• " + b
                    p.font.size = Pt(12)
        else:
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = s["title"]
            body = slide.placeholders[1].text_frame
            body.word_wrap = True
            for i, b in enumerate(s.get("bullets", []) or ["(no facts)"]):
                p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                p.text = b
                p.font.size = Pt(12)
    prs.save(out_path)
    return True


# ----- main -----------------------------------------------------------------

def main(argv):
    ap = argparse.ArgumentParser(description="A2: findings.json -> deck")
    ap.add_argument("findings", help="path to findings.json")
    ap.add_argument("--emit",
                    choices=["auto", "payload", "pptx", "json", "emit-args"],
                    default="auto")
    ap.add_argument("--out", default=None, help="output path")
    ap.add_argument("--title", default=None, help="deck title override")
    ap.add_argument("--company", default=None, help="company name for subtitle")
    # A5 deck-OUT: --emit emit-args builds the exact emit_artifact arguments
    # body from an already-rendered .pptx (from rockie-pptagent-render).
    ap.add_argument("--pptx-in", default=None,
                    help="path to a rendered .pptx (required for --emit emit-args)")
    ap.add_argument("--notebook-id",
                    default=os.environ.get("PLATFORM_LAB_ID")
                    or os.environ.get("LAB_ID"),
                    help="lab/notebook id for emit_artifact (default: "
                         "$PLATFORM_LAB_ID / $LAB_ID)")
    args = ap.parse_args(argv)

    findings, by_id = load_findings(args.findings)
    title = args.title or (by_id[0].get("headline")
                           or "Acquisition due-diligence findings")
    slides = build_slides(findings, by_id, title, args.company)
    n = len(slides)

    if n > HARD_SLIDE_CAP:
        die("deck has %d slides; powerpoint skill caps at %d" % (n, HARD_SLIDE_CAP))
    if n < TARGET_MIN or n > TARGET_MAX:
        sys.stderr.write(
            "render_deck: NOTE: %d slides (A2 target band is %d-%d)\n"
            % (n, TARGET_MIN, TARGET_MAX))

    # assertions that protect the partner-grade contract
    assert slides[0]["section_id"] == 0, "Section 0 must lead the deck"

    emit = args.emit
    if emit == "auto":
        try:
            import pptx  # noqa: F401
            emit = "pptx"
        except Exception:
            emit = "json"

    if emit == "emit-args":
        # A5 deck-OUT: produce the ready-to-post emit_artifact arguments from
        # an already-rendered .pptx. This is the OUT half of the connector --
        # the agent posts this body to /api/agent-tools/emit_artifact (or calls
        # the emit_artifact MCP tool) and the deck ships to the user.
        if not args.pptx_in:
            die("--emit emit-args requires --pptx-in <rendered .pptx>")
        if not os.path.isfile(args.pptx_in):
            die("--pptx-in not found: %s" % args.pptx_in)
        if not args.notebook_id:
            die("emit_artifact needs a notebook_id; set --notebook-id or "
                "$PLATFORM_LAB_ID in the runtime")
        # emit_artifact rejects names not matching ^[A-Za-z0-9][A-Za-z0-9._ -]{0,199}$
        # (rockie-pptagent temp names can lead with _ or carry unsafe chars), so
        # sanitize the basename before emitting or the OUT half hard-fails at A6.
        _alnum = "abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789"
        _raw = os.path.basename(args.pptx_in)
        fname = "".join(c if c in (_alnum + "._ -") else "-" for c in _raw)
        if not fname or fname[0] not in _alnum:
            fname = "deck-" + fname
        if not fname.lower().endswith(".pptx"):
            fname = fname[:195] + ".pptx"
        fname = fname[:200]
        if not fname.lower().endswith(".pptx"):
            fname = fname[:195] + ".pptx"
        emit_args = build_emit_args(args.pptx_in, title, n,
                                    args.notebook_id, filename=fname)
        out = args.out or default_out(".json")
        with open(out, "w", encoding="utf-8") as fh:
            json.dump({"tool": "emit_artifact", "arguments": emit_args},
                      fh, indent=2, ensure_ascii=False)
        print("render_deck: emit_artifact args -> %s (%d slides, lab=%s)"
              % (out, n, args.notebook_id))
        print("SLIDES=%d SECTION0_LEADS=%s EMIT=%s"
              % (n, str(slides[0]["section_id"] == 0).lower(), emit))
        return 0

    if emit == "payload":
        payload = to_pptagent_payload(slides, title)
        out = args.out or default_out(".json")
        with open(out, "w", encoding="utf-8") as fh:
            json.dump(payload, fh, indent=2, ensure_ascii=False)
        print("render_deck: powerpoint-skill payload -> %s (%d slides)" % (out, n))
    elif emit == "pptx":
        out = args.out or default_out(".pptx")
        if not render_pptx(slides, out, title):
            sys.stderr.write("render_deck: python-pptx unavailable; use --emit json\n")
            sys.exit(3)
        print("render_deck: rendered .pptx -> %s (%d slides)" % (out, n))
    else:  # json
        out = args.out or default_out(".json")
        with open(out, "w", encoding="utf-8") as fh:
            json.dump({"title": title, "slide_count": n, "slides": slides},
                      fh, indent=2, ensure_ascii=False)
        print("render_deck: deck spec -> %s (%d slides)" % (out, n))

    # machine-checkable summary line for CI / verification
    print("SLIDES=%d SECTION0_LEADS=%s EMIT=%s"
          % (n, str(slides[0]["section_id"] == 0).lower(), emit))
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
